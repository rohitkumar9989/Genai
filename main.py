import streamlit as st
import pandas as pd
import numpy
from initialize_ import *
import matplotlib.pyplot as plt
import seaborn as sns
import re
import numpy as np
import time
import io
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.util import Pt
from langchain_core.output_parsers import StrOutputParser

def main():

    st.title("Effective Data Visualization Tool")
    if os.path.exists("datasets")==False:
        os.mkdir("datasets")
    if 'llm_model1' not in st.session_state:
        llm = llm_create()
        with st.spinner("Loading the model"):
            st.session_state.llm_model1 = llm.subroutine_main()
            st.session_state.questionare_lw_model = llm.questionare_lightweight_model()
    else:
        llm_model1 = st.session_state.llm_model1
        questionare_lw_model = st.session_state.questionare_lw_model
    create_template=Create_prompt_template()
    csv_data=st.file_uploader(label="Enter the csv file which you want to create visualization on: ")
    prs = Presentation()
    save_ppt=st.radio("Create a presentation?: ", ["Yes", "No"], captions=["Dont worry we create effective Viz", "ok"])

    if csv_data:
        file_name=csv_data.name
        flag_ste=True

        first_iteration=True
        if "state" not in st.session_state:
            st.session_state.state=file_name
        else:
            first_iteration=False
            if st.session_state["state"]==file_name:
                flag_ste=False
            else:
                st.session_state.state=file_name
                flag_ste=True
        st.write(file_name, st.session_state["state"])
        file_path=os.path.join("datasets", csv_data.name)
        if os.path.exists(file_path)==False:
            st.write("Please move the csv data into the `datasets` folder")
        else:
            data=pd.read_csv(file_path)
            pandas=Process_pandas(data)
            flag, null=pandas.check_nan()
            if flag:
                message=st.warning("Nan Values are present. Performing Processing.")
                data=pandas.preprocess_data(null)
                time.sleep(3)
                message.success("Done Processing")
            questioner=Ask_questions(file_path)
            ppt=Makeppt(file_path)
            class Visualize_data():
                def run(self, data):
                    for key in data:
                        try:
                            sub_dictionary = data[key]
                            xaxis = sub_dictionary.get("X-Axis")
                            yaxis = sub_dictionary.get("Y-Axis")
                            xaxis_values, yaxis_values = pandas.obtain_data(xaxis, yaxis)
                            viz_type = sub_dictionary.get("Type")
                            desc=sub_dictionary.get("desc")
                            if viz_type == "Scatter Plot":
                                data_df = pd.DataFrame({xaxis: xaxis_values, yaxis: yaxis_values})
                                aggregated_data = data_df.groupby(xaxis).sum().reset_index()
                                aggregated_data = aggregated_data.sort_values(by=yaxis, ascending=False)
                                aggregated_data = aggregated_data.head(50)
                                xaxis_values = aggregated_data[xaxis]
                                yaxis_values = aggregated_data[yaxis]
                                fig, ax = plt.subplots(figsize=(12, 6))
                                sns.scatterplot(x=xaxis_values, y=yaxis_values, color='blue', alpha=0.6, edgecolor='black', ax=ax)
                                ax.set_title("Scatter Plot")
                                ax.set_xlabel(xaxis)
                                ax.set_ylabel(yaxis)
                                st.subheader(desc)
                                plt.xticks(rotation=90)
                                st.pyplot(fig)
                            elif viz_type == "Bar Chart":
                                data_df = pd.DataFrame({xaxis: xaxis_values, yaxis: yaxis_values})
            
                                if re.search(r"(.*?)-(.*?)", str(xaxis_values)) or re.search(r"(.*?)/(.*?)", str(xaxis_values)):
                                    try:
                                        data_df[xaxis] = pd.to_datetime(data_df[xaxis], errors='coerce')
                                        df1 = data_df.copy()
                                        df1['group_key'] = df1[xaxis].dt.year
                                        aggregated_data = df1.groupby('group_key')[yaxis].sum().reset_index()
                                    except Exception as e:
                                        st.error(f"Error converting {xaxis} to datetime: {e}")
                                        aggregated_data = data_df.groupby(xaxis).sum().reset_index()
                                else:
                                    aggregated_data = data_df.groupby(xaxis).sum().reset_index()
                                aggregated_data = aggregated_data.sort_values(by=yaxis, ascending=False)
                                aggregated_data = aggregated_data.head(50)
                                xaxis_values = aggregated_data.iloc[:, 0]
                                yaxis_values = aggregated_data[yaxis]
                                fig, ax = plt.subplots(figsize=(12, 6))
                                sns.barplot(x=xaxis_values, y=yaxis_values, color='blue', ax=ax)
                                ax.set_xlabel(xaxis)
                                ax.set_ylabel(yaxis)
                                ax.set_title("Bar Chart Example")
                                st.subheader(desc)
                                plt.xticks(rotation=90)
                                st.pyplot(fig)
                            elif viz_type == "Pie Chart":
                                data_df = pd.DataFrame({xaxis: xaxis_values, yaxis: yaxis_values})
                                aggregated_data = data_df.groupby(xaxis).sum().reset_index()
                                pie_data = aggregated_data[:10]
                                labels = pie_data[xaxis]
                                sizes = pie_data[yaxis]
                                fig, ax = plt.subplots()
                                wedges, texts, autotexts = ax.pie(
                                    sizes,
                                    labels=labels,
                                    autopct='%1.1f%%',
                                    startangle=90,
                                    colors=sns.color_palette("pastel"),
                                    textprops=dict(color="black"),
                                )
                                ax.legend(
                                    wedges,
                                    labels,
                                    title=xaxis,
                                    loc="center left",
                                    bbox_to_anchor=(1, 0.5),
                                )
                                ax.set_title(f"Pie Chart of {xaxis} vs {yaxis}")
                                st.subheader(desc)
                                st.pyplot(fig)
                            elif viz_type == "Line Chart":
                                data_df = pd.DataFrame({xaxis: xaxis_values, yaxis: yaxis_values})
                                if re.search(r"(.*?)-(.*?)", str(xaxis_values)) or re.search(r"(.*?)/(.*?)", str(xaxis_values)):
                                    try:
                                        data_df[xaxis] = pd.to_datetime(data_df[xaxis], errors='coerce')
                                        df1 = data_df.copy()
                                        df1['group_key'] = df1[xaxis].dt.year
                                        aggregated_data = df1.groupby('group_key')[yaxis].sum().reset_index()
                                    except Exception as e:
                                        st.error(f"Error converting {xaxis} to datetime: {e}")
                                        aggregated_data = data_df.groupby(xaxis).sum().reset_index()
                                else:
                                    aggregated_data = data_df.groupby(xaxis).sum().reset_index()
                                aggregated_data = data_df.groupby(xaxis).sum().reset_index()
                                aggregated_data = aggregated_data.sort_values(by=yaxis, ascending=False)
                                aggregated_data = aggregated_data.head(50)
                                xaxis_values = aggregated_data.iloc[:, 0]
                                yaxis_values = aggregated_data[yaxis]
                                fig, ax = plt.subplots(figsize=(15, 8))
                                sns.lineplot(x=xaxis_values, y=yaxis_values, marker="o", ax=ax)
                                ax.set_title("Line Chart")
                                ax.set_xlabel(xaxis)
                                ax.set_ylabel(yaxis)
                                plt.xticks(rotation=90)
                                st.subheader(desc)
                                st.pyplot(fig)
                            else:
                                continue
                            buffer = io.BytesIO()
                            plt.savefig(buffer, format="png")
                            buffer.seek(0)

                            st.download_button(
                                label="Download Report",
                                data=buffer,
                                file_name=f"{viz_type.replace(' ', '_')}_{xaxis}_vs_{yaxis}.png",
                                mime="image/png"
                            )
                            if save_ppt == "Yes":
                                pptx_image = f"{viz_type.replace(' ', '_')}_{xaxis}_vs_{yaxis}.png"
                                with open(pptx_image, "wb") as f:
                                    f.write(buffer.getvalue())
                                slide_layout = prs.slide_layouts[5]  
                                slide = prs.slides.add_slide(slide_layout)
                                shapes = slide.shapes

                                left = Inches(0.5)
                                top = Inches(3)
                                pic = slide.shapes.add_picture(pptx_image, left, top, width=Inches(4.5))

                                text_box = shapes.add_textbox(Inches(5.5), Inches(1), Inches(4), Inches(3))  
                                text_frame = text_box.text_frame
                                answer = ppt.presentation(desc)
                                text_frame.text = f"{viz_type} of {xaxis} vs {yaxis}\n\n{answer}"

                                
                                text_frame.word_wrap = True
                                text_frame.auto_size = True

                                
                                for paragraph in text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.size = Pt(12)
                                        run.font.color.rgb = RGBColor(0, 0, 0)  

                                os.remove(pptx_image)

                        except Exception as e:
                            print (e)
                    if save_ppt=="Yes":
                        pptx_path = "data_visualization_presentation.pptx"
                        prs.save(pptx_path)
                        with open(pptx_path, "rb") as f:
                            st.download_button(
                                label="Download Presentation",
                                data=f,
                                file_name="Data_Visualization.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            )
            visualizer=Visualize_data()
            if flag_ste or first_iteration:
                dict_=pandas.get_dictionary()
                promt_template=create_template.call(dict_)
                chain=promt_template|llm_model1
                dictionary=", ".join(f"Column Name: {key} Its datatype: {value} \n" for key, value in dict_.items())
                with st.spinner("Thinking..."):
                    output=chain.invoke({"context": dictionary, "values":data.iloc[:50]})
                    pars=StrOutputParser()
                    output=pars.invoke(output)
                    extractor=Extract_data(output)
                    st.session_state.viz_dict=extractor.extract_columns()
                    with st.spinner("Visualizing the data"):
                        visualizer.run(st.session_state.viz_dict)
            else:
                with st.spinner("Visualizing the data"):
                        visualizer.run(st.session_state.viz_dict)


            
            col1, col2 = st.columns([4, 1])  

            with col1:
                text = st.text_input("Enter your Question")

            with col2:
                model_choice = st.selectbox("Model", ["ChatGPT", "Gemma", "Google Tapas"])

            if text and model_choice:
                question = questioner.answer_question(text, model_choice)
                st.write(f"Using {model_choice}:")
                st.write(question)
    else:
        pass




